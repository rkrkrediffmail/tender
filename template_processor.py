# template_processor.py
"""
Template processing engine for DOCX and PPTX proposal generation with bookmark replacement
"""

import os
import logging
import hashlib
from typing import Dict, List, Any, Optional, Tuple
from datetime import datetime
import zipfile
import xml.etree.ElementTree as ET
from xml.dom import minidom
import re

# Import with fallback for different python-docx versions
try:
    from docx import Document
    from docx.shared import Inches, Pt
    from docx.enum.text import WD_PARAGRAPH_ALIGNMENT
    from docx.oxml.shared import qn
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False
    logging.warning("python-docx not available - DOCX template processing disabled")

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    logging.warning("python-pptx not available - PPTX template processing disabled")

from models import db, ProposalTemplate, TemplateBookmark

logger = logging.getLogger(__name__)

class TemplateProcessor:
    """
    Advanced template processor for DOCX and PPTX files with bookmark replacement
    """
    
    def __init__(self):
        self.supported_formats = []
        if DOCX_AVAILABLE:
            self.supported_formats.append('docx')
        if PPTX_AVAILABLE:
            self.supported_formats.append('pptx')
            
        logger.info(f"Template processor initialized with support for: {', '.join(self.supported_formats)}")
    
    def process_template_upload(self, 
                              file_path: str,
                              original_filename: str,
                              template_info: Dict[str, Any],
                              user_id: int) -> Dict[str, Any]:
        """
        Process uploaded template file and extract bookmarks
        
        Args:
            file_path: Path to uploaded template file
            original_filename: Original filename
            template_info: Template metadata
            user_id: User who uploaded the template
            
        Returns:
            Dict with processing results
        """
        try:
            # Determine file type
            file_extension = original_filename.lower().split('.')[-1]
            if file_extension not in self.supported_formats:
                return {
                    'success': False,
                    'error': f'Unsupported file format: {file_extension}. Supported formats: {", ".join(self.supported_formats)}'
                }
            
            # Calculate file hash
            file_hash = self._calculate_file_hash(file_path)
            file_size = os.path.getsize(file_path)
            
            # Extract bookmarks from template
            if file_extension == 'docx':
                bookmarks = self._extract_docx_bookmarks(file_path)
            elif file_extension == 'pptx':
                bookmarks = self._extract_pptx_placeholders(file_path)
            else:
                bookmarks = []
            
            # Create template record
            template = ProposalTemplate(
                name=template_info.get('name', original_filename),
                description=template_info.get('description', ''),
                template_type=file_extension,
                category=template_info.get('category', 'general'),
                filename=os.path.basename(file_path),
                original_filename=original_filename,
                file_path=file_path,
                file_size=file_size,
                file_hash=file_hash,
                is_active=template_info.get('is_active', True),
                is_default=template_info.get('is_default', False),
                company_info=template_info.get('company_info', {}),
                uploaded_by=user_id
            )
            
            db.session.add(template)
            db.session.flush()  # Get template ID
            
            # Create bookmark records
            for bookmark_info in bookmarks:
                bookmark = TemplateBookmark(
                    template_id=template.id,
                    bookmark_name=bookmark_info['name'],
                    display_name=bookmark_info.get('display_name', bookmark_info['name']),
                    description=bookmark_info.get('description', f'Content for {bookmark_info["name"]}'),
                    content_type=bookmark_info.get('content_type', 'dynamic'),
                    content_source=bookmark_info.get('content_source', 'ai_generated'),
                    content_format=bookmark_info.get('content_format', 'text'),
                    is_required=bookmark_info.get('is_required', True),
                    default_content=bookmark_info.get('default_content', ''),
                    ai_prompt_template=bookmark_info.get('ai_prompt_template', '')
                )
                db.session.add(bookmark)
            
            db.session.commit()
            
            result = {
                'success': True,
                'template_id': template.template_id,
                'template_db_id': template.id,
                'bookmarks_found': len(bookmarks),
                'bookmarks': [b['name'] for b in bookmarks],
                'file_size': file_size,
                'file_hash': file_hash,
                'message': f'Template uploaded successfully with {len(bookmarks)} bookmarks detected'
            }
            
            logger.info(f"Template processed successfully: {original_filename} - {result['message']}")
            return result
            
        except Exception as e:
            db.session.rollback()
            logger.error(f"Error processing template {original_filename}: {e}")
            return {
                'success': False,
                'error': str(e),
                'message': f'Failed to process template: {str(e)}'
            }
    
    def generate_proposal_from_template(self,
                                      template_id: int,
                                      project_data: Dict[str, Any],
                                      content_data: Dict[str, Any],
                                      output_path: str) -> Dict[str, Any]:
        """
        Generate proposal by filling template bookmarks with content
        
        Args:
            template_id: Database ID of the template
            project_data: Project information
            content_data: Content to fill into bookmarks
            output_path: Path where to save the generated file
            
        Returns:
            Dict with generation results
        """
        try:
            # Load template from database
            template = ProposalTemplate.query.get(template_id)
            if not template:
                return {'success': False, 'error': 'Template not found'}
            
            if not os.path.exists(template.file_path):
                return {'success': False, 'error': 'Template file not found on disk'}
            
            # Process based on template type
            if template.template_type == 'docx':
                result = self._process_docx_template(template, project_data, content_data, output_path)
            elif template.template_type == 'pptx':
                result = self._process_pptx_template(template, project_data, content_data, output_path)
            else:
                return {'success': False, 'error': f'Unsupported template type: {template.template_type}'}
            
            # Update template usage statistics
            if result.get('success'):
                template.usage_count += 1
                template.last_used_at = datetime.utcnow()
                db.session.commit()
            
            return result
            
        except Exception as e:
            logger.error(f"Error generating proposal from template {template_id}: {e}")
            return {
                'success': False,
                'error': str(e),
                'message': f'Failed to generate proposal: {str(e)}'
            }
    
    def _extract_docx_bookmarks(self, file_path: str) -> List[Dict[str, Any]]:
        """Extract bookmarks from DOCX file"""
        bookmarks = []
        
        if not DOCX_AVAILABLE:
            logger.warning("python-docx not available, cannot extract DOCX bookmarks")
            return bookmarks
        
        try:
            # Open as ZIP file to access XML directly
            with zipfile.ZipFile(file_path, 'r') as zip_file:
                # Read document.xml
                doc_xml = zip_file.read('word/document.xml').decode('utf-8')
                
                # Parse XML to find bookmarks
                root = ET.fromstring(doc_xml)
                
                # Define namespaces
                namespaces = {
                    'w': 'http://schemas.openxmlformats.org/wordprocessingml/2006/main'
                }
                
                # Find bookmark starts
                bookmark_starts = root.findall('.//w:bookmarkStart', namespaces)
                
                for bookmark_start in bookmark_starts:
                    name = bookmark_start.get('{http://schemas.openxmlformats.org/wordprocessingml/2006/main}name')
                    if name and not name.startswith('_'):  # Skip internal bookmarks
                        bookmarks.append({
                            'name': name,
                            'display_name': name.replace('_', ' ').title(),
                            'content_type': 'dynamic',
                            'content_source': self._infer_content_source(name),
                            'content_format': 'text',
                            'is_required': True
                        })
                
                logger.info(f"Extracted {len(bookmarks)} bookmarks from DOCX: {[b['name'] for b in bookmarks]}")
                
        except Exception as e:
            logger.error(f"Error extracting DOCX bookmarks: {e}")
            # Fallback: try with python-docx
            try:
                doc = Document(file_path)
                # This is a simplified approach - python-docx doesn't have great bookmark support
                # We'll look for text that looks like bookmark placeholders
                for paragraph in doc.paragraphs:
                    text = paragraph.text
                    # Look for {{bookmark}} patterns
                    matches = re.findall(r'\{\{([^}]+)\}\}', text)
                    for match in matches:
                        bookmark_name = match.strip()
                        if bookmark_name not in [b['name'] for b in bookmarks]:
                            bookmarks.append({
                                'name': bookmark_name,
                                'display_name': bookmark_name.replace('_', ' ').title(),
                                'content_type': 'dynamic',
                                'content_source': self._infer_content_source(bookmark_name),
                                'content_format': 'text',
                                'is_required': True
                            })
                
            except Exception as fallback_error:
                logger.error(f"Fallback bookmark extraction also failed: {fallback_error}")
        
        return bookmarks
    
    def _extract_pptx_placeholders(self, file_path: str) -> List[Dict[str, Any]]:
        """Extract placeholders from PPTX file"""
        placeholders = []
        
        if not PPTX_AVAILABLE:
            logger.warning("python-pptx not available, cannot extract PPTX placeholders")
            return placeholders
        
        try:
            prs = Presentation(file_path)
            
            # Track found placeholders to avoid duplicates
            found_placeholders = set()
            
            # Iterate through slides
            for slide_idx, slide in enumerate(prs.slides):
                # Check shapes in slide
                for shape in slide.shapes:
                    if hasattr(shape, 'text'):
                        text = shape.text
                        # Look for {{placeholder}} patterns
                        matches = re.findall(r'\{\{([^}]+)\}\}', text)
                        for match in matches:
                            placeholder_name = match.strip()
                            if placeholder_name not in found_placeholders:
                                found_placeholders.add(placeholder_name)
                                placeholders.append({
                                    'name': placeholder_name,
                                    'display_name': placeholder_name.replace('_', ' ').title(),
                                    'content_type': 'dynamic',
                                    'content_source': self._infer_content_source(placeholder_name),
                                    'content_format': 'text',
                                    'is_required': True,
                                    'location': f'Slide {slide_idx + 1}'
                                })
            
            logger.info(f"Extracted {len(placeholders)} placeholders from PPTX: {[p['name'] for p in placeholders]}")
            
        except Exception as e:
            logger.error(f"Error extracting PPTX placeholders: {e}")
        
        return placeholders
    
    def _process_docx_template(self,
                             template: ProposalTemplate,
                             project_data: Dict[str, Any],
                             content_data: Dict[str, Any],
                             output_path: str) -> Dict[str, Any]:
        """Process DOCX template with content replacement"""
        
        if not DOCX_AVAILABLE:
            return {'success': False, 'error': 'python-docx not available'}
        
        try:
            # Load the template document
            doc = Document(template.file_path)
            
            # Get bookmarks for this template
            bookmarks = TemplateBookmark.query.filter_by(template_id=template.id).all()
            bookmark_content = {}
            
            # Process each bookmark
            for bookmark in bookmarks:
                content = self._generate_bookmark_content(bookmark, project_data, content_data)
                bookmark_content[bookmark.bookmark_name] = content
                
                # Replace content in document
                self._replace_docx_content(doc, bookmark.bookmark_name, content)
            
            # Also replace {{placeholder}} patterns for compatibility
            for paragraph in doc.paragraphs:
                for bookmark_name, content in bookmark_content.items():
                    placeholder = f'{{{{{bookmark_name}}}}}'
                    if placeholder in paragraph.text:
                        paragraph.text = paragraph.text.replace(placeholder, str(content))
            
            # Replace in tables
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        for paragraph in cell.paragraphs:
                            for bookmark_name, content in bookmark_content.items():
                                placeholder = f'{{{{{bookmark_name}}}}}'
                                if placeholder in paragraph.text:
                                    paragraph.text = paragraph.text.replace(placeholder, str(content))
            
            # Save the document
            doc.save(output_path)
            
            return {
                'success': True,
                'output_path': output_path,
                'bookmarks_processed': len(bookmark_content),
                'bookmark_content': bookmark_content,
                'file_size': os.path.getsize(output_path)
            }
            
        except Exception as e:
            logger.error(f"Error processing DOCX template: {e}")
            return {'success': False, 'error': str(e)}
    
    def _process_pptx_template(self,
                             template: ProposalTemplate,
                             project_data: Dict[str, Any],
                             content_data: Dict[str, Any],
                             output_path: str) -> Dict[str, Any]:
        """Process PPTX template with content replacement"""
        
        if not PPTX_AVAILABLE:
            return {'success': False, 'error': 'python-pptx not available'}
        
        try:
            # Load the template presentation
            prs = Presentation(template.file_path)
            
            # Get bookmarks for this template
            bookmarks = TemplateBookmark.query.filter_by(template_id=template.id).all()
            bookmark_content = {}
            
            # Process each bookmark
            for bookmark in bookmarks:
                content = self._generate_bookmark_content(bookmark, project_data, content_data)
                bookmark_content[bookmark.bookmark_name] = content
            
            # Replace content in slides
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, 'text'):
                        original_text = shape.text
                        updated_text = original_text
                        
                        # Replace all placeholders in this shape
                        for bookmark_name, content in bookmark_content.items():
                            placeholder = f'{{{{{bookmark_name}}}}}'
                            if placeholder in updated_text:
                                updated_text = updated_text.replace(placeholder, str(content))
                        
                        # Update the shape text if it changed
                        if updated_text != original_text:
                            shape.text = updated_text
                    
                    # Handle text in tables within slides
                    if shape.has_table:
                        for row in shape.table.rows:
                            for cell in row.cells:
                                original_text = cell.text
                                updated_text = original_text
                                
                                for bookmark_name, content in bookmark_content.items():
                                    placeholder = f'{{{{{bookmark_name}}}}}'
                                    if placeholder in updated_text:
                                        updated_text = updated_text.replace(placeholder, str(content))
                                
                                if updated_text != original_text:
                                    cell.text = updated_text
            
            # Save the presentation
            prs.save(output_path)
            
            return {
                'success': True,
                'output_path': output_path,
                'bookmarks_processed': len(bookmark_content),
                'bookmark_content': bookmark_content,
                'file_size': os.path.getsize(output_path)
            }
            
        except Exception as e:
            logger.error(f"Error processing PPTX template: {e}")
            return {'success': False, 'error': str(e)}
    
    def _replace_docx_content(self, doc: 'Document', bookmark_name: str, content: str):
        """Replace content at bookmark location in DOCX (simplified approach)"""
        # This is a simplified implementation
        # A full implementation would require deeper XML manipulation
        placeholder = f'{{{{{bookmark_name}}}}}'
        
        # Replace in paragraphs
        for paragraph in doc.paragraphs:
            if placeholder in paragraph.text:
                paragraph.text = paragraph.text.replace(placeholder, content)
    
    def _generate_bookmark_content(self,
                                 bookmark: TemplateBookmark,
                                 project_data: Dict[str, Any],
                                 content_data: Dict[str, Any]) -> str:
        """Generate content for a specific bookmark"""
        try:
            # Check if content is provided in content_data
            if bookmark.bookmark_name in content_data:
                content = content_data[bookmark.bookmark_name]
                if content:
                    return str(content)
            
            # Generate content based on content_source
            if bookmark.content_type == 'static':
                return bookmark.default_content or ''
            
            elif bookmark.content_type == 'dynamic':
                return self._generate_dynamic_content(bookmark, project_data, content_data)
            
            elif bookmark.content_type == 'ai_generated':
                return self._generate_ai_content(bookmark, project_data, content_data)
            
            else:
                return bookmark.default_content or f'[{bookmark.display_name}]'
                
        except Exception as e:
            logger.error(f"Error generating content for bookmark {bookmark.bookmark_name}: {e}")
            return bookmark.default_content or f'[{bookmark.display_name}]'
    
    def _generate_dynamic_content(self,
                                bookmark: TemplateBookmark,
                                project_data: Dict[str, Any],
                                content_data: Dict[str, Any]) -> str:
        """Generate dynamic content based on project data"""
        
        source = bookmark.content_source
        
        # Map content sources to project data
        content_mapping = {
            'project_name': project_data.get('name', 'Unnamed Project'),
            'project_description': project_data.get('description', ''),
            'client_name': project_data.get('client_name', 'Valued Client'),
            'company_name': content_data.get('company_name', 'Your Company'),
            'contact_person': content_data.get('contact_person', 'Project Manager'),
            'project_date': datetime.now().strftime('%B %d, %Y'),
            'current_date': datetime.now().strftime('%B %d, %Y'),
            'current_year': str(datetime.now().year)
        }
        
        # Check for company info from template
        if hasattr(bookmark, 'template') and bookmark.template.company_info:
            company_info = bookmark.template.company_info
            content_mapping.update({
                'company_address': company_info.get('address', ''),
                'company_phone': company_info.get('phone', ''),
                'company_email': company_info.get('email', ''),
                'company_website': company_info.get('website', '')
            })
        
        return content_mapping.get(source, bookmark.default_content or f'[{source}]')
    
    def _generate_ai_content(self,
                           bookmark: TemplateBookmark,
                           project_data: Dict[str, Any],
                           content_data: Dict[str, Any]) -> str:
        """Generate AI content for bookmark"""
        
        # If AI content is provided in content_data, use it
        ai_key = f'{bookmark.bookmark_name}_ai_content'
        if ai_key in content_data:
            return content_data[ai_key]
        
        # For now, return placeholder - this would integrate with AI generation
        # In a full implementation, this would call Claude API with the bookmark's ai_prompt_template
        return bookmark.default_content or f'[AI-generated content for {bookmark.display_name}]'
    
    def _infer_content_source(self, bookmark_name: str) -> str:
        """Infer likely content source from bookmark name"""
        name_lower = bookmark_name.lower()
        
        # Common patterns
        if any(word in name_lower for word in ['project', 'title', 'name']):
            return 'project_name'
        elif any(word in name_lower for word in ['client', 'customer']):
            return 'client_name'
        elif any(word in name_lower for word in ['company', 'organization', 'vendor']):
            return 'company_name'
        elif any(word in name_lower for word in ['contact', 'person', 'manager']):
            return 'contact_person'
        elif any(word in name_lower for word in ['date', 'today']):
            return 'current_date'
        elif any(word in name_lower for word in ['description', 'summary', 'overview']):
            return 'project_analysis'
        elif any(word in name_lower for word in ['requirement', 'spec', 'need']):
            return 'requirements'
        elif any(word in name_lower for word in ['solution', 'approach', 'method']):
            return 'ai_generated'
        else:
            return 'ai_generated'
    
    def _calculate_file_hash(self, file_path: str) -> str:
        """Calculate SHA-256 hash of file"""
        try:
            sha256_hash = hashlib.sha256()
            with open(file_path, "rb") as f:
                for byte_block in iter(lambda: f.read(4096), b""):
                    sha256_hash.update(byte_block)
            return sha256_hash.hexdigest()
        except Exception as e:
            logger.error(f"Error calculating file hash: {e}")
            return ""
    
    def validate_template(self, file_path: str) -> Dict[str, Any]:
        """Validate uploaded template file"""
        try:
            if not os.path.exists(file_path):
                return {'valid': False, 'error': 'File does not exist'}
            
            file_extension = file_path.lower().split('.')[-1]
            if file_extension not in self.supported_formats:
                return {
                    'valid': False,
                    'error': f'Unsupported format: {file_extension}. Supported: {", ".join(self.supported_formats)}'
                }
            
            # Try to open the file to validate it
            if file_extension == 'docx' and DOCX_AVAILABLE:
                try:
                    doc = Document(file_path)
                    # Basic validation - check if we can access paragraphs
                    paragraph_count = len(doc.paragraphs)
                    return {
                        'valid': True,
                        'file_type': 'docx',
                        'paragraph_count': paragraph_count,
                        'has_tables': len(doc.tables) > 0
                    }
                except Exception as e:
                    return {'valid': False, 'error': f'Invalid DOCX file: {str(e)}'}
            
            elif file_extension == 'pptx' and PPTX_AVAILABLE:
                try:
                    prs = Presentation(file_path)
                    # Basic validation - check slides
                    slide_count = len(prs.slides)
                    return {
                        'valid': True,
                        'file_type': 'pptx',
                        'slide_count': slide_count
                    }
                except Exception as e:
                    return {'valid': False, 'error': f'Invalid PPTX file: {str(e)}'}
            
            else:
                return {'valid': False, 'error': f'No processor available for {file_extension}'}
                
        except Exception as e:
            return {'valid': False, 'error': str(e)}

# Singleton instance
template_processor_instance = None

def get_template_processor() -> TemplateProcessor:
    """Get or create template processor instance"""
    global template_processor_instance
    if template_processor_instance is None:
        template_processor_instance = TemplateProcessor()
    return template_processor_instance